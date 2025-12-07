from pathlib import Path
import os
import re
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE_TYPE

from ..utils.state import State
from ..utils.slides_as_png import export_slide_as_png

def node_parse_all(state: State) -> State:

    
    pptx_path = state["pptx_path"]
    work_dir = state.get("work_dir", "./")

    # === 0) 폴더 생성 ===
    base_dir = Path(work_dir).expanduser().resolve()
    slides_dir = base_dir / "slides"
    media_dir = base_dir / "media"
    slides_dir.mkdir(parents=True, exist_ok=True)
    media_dir.mkdir(parents=True, exist_ok=True)

    globals()["SLIDES_DIR"] = str(slides_dir)
    globals()["MEDIA_DIR"] = str(media_dir)

    prs = Presentation(pptx_path)
    slide_count = len(prs.slides)

    print(f"총 {slide_count}장의 슬라이드를 처리합니다.")

    # ========= 평탄화 리스트 =========
    titles_list = []
    texts_list = []
    tables_list = []
    images_list = []
    snapshots_list = []
    shape_texts_list = []
    links_list = []

    url_pattern = r"https?://[^\s]+"

    # ---------------------------------------------
    # 슬라이드 단위 반복
    # ---------------------------------------------
    for slide_idx, slide in enumerate(prs.slides):
        print(f"\n=== {slide_idx+1}번째 슬라이드 처리 중 ===")

        # 1) PNG 생성
        slide_state = {
            "pptx_path": pptx_path,
            "work_dir": str(slides_dir),
            "slide_index": slide_idx,
        }
        slide_state = export_slide_as_png(slide_state)

        src_png = slide_state["slide_image"]
        dst_png = slides_dir / f"slide_{slide_idx+1}.png"

        if os.path.exists(src_png):
            os.replace(src_png, dst_png)
            snapshot_path = str(dst_png)
        else:
            snapshot_path = ""

        # ============================
        # 요소 수집
        # ============================
        slide_title = ""
        body_texts = []
        shape_texts = []
        slide_tables = []
        slide_images = []
        slide_links = set()

        # ---- 제목 추출 ----
        for sh in slide.shapes:
            try:
                if sh.is_placeholder and sh.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                    if sh.has_text_frame:
                        t = sh.text.strip()
                        if t:
                            slide_title = t
            except:
                pass

        # ---- 도형 텍스트 ----
        def collect_shape_texts(shape):
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for sub in shape.shapes:
                    collect_shape_texts(sub)
                return

            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    shape_texts.append(t)

        for sh in slide.shapes:
            collect_shape_texts(sh)

        # ---- 링크 ----
        for sh in slide.shapes:
            if sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    runs_text = "".join(r.text for r in p.runs)
                    found = re.findall(url_pattern, runs_text)
                    if found:
                        slide_links.update(found)

            # 도형 클릭 하이퍼링크
            try:
                if hasattr(sh, "click_action") and sh.click_action.hyperlink.address:
                    slide_links.add(sh.click_action.hyperlink.address)
            except:
                pass

        # ---- 본문 텍스트(제목/도형/링크 제외) ----
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue

            # 제목 제외
            try:
                if sh.is_placeholder and sh.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                    continue
            except:
                pass

            for p in sh.text_frame.paragraphs:
                txt = "".join(r.text for r in p.runs).strip()

                if not txt:
                    continue
                if txt in shape_texts:
                    continue
                if re.search(url_pattern, txt):
                    continue

                body_texts.append(txt)

        # ---- 표 ----
        for sh in slide.shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.TABLE:
                tbl = [[cell.text.strip() for cell in row.cells] for row in sh.table.rows]
                slide_tables.append(tbl)

        # ---- 이미지 ----
        for sh in slide.shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    ext = sh.image.ext
                except:
                    ext = "png"

                img_name = f"slide{slide_idx+1}_img_{len(slide_images)+1}.{ext}"
                img_path = media_dir / img_name
                with open(img_path, "wb") as f:
                    f.write(sh.image.blob)
                slide_images.append(str(img_path))

        # ---- 평탄화 리스트에 저장 ----
        titles_list.append(slide_title)
        texts_list.append("\n".join(body_texts).strip())
        tables_list.append(slide_tables)
        images_list.append(slide_images)
        snapshots_list.append(snapshot_path)
        shape_texts_list.append(list(dict.fromkeys(shape_texts)))
        links_list.append(list(slide_links))

        print(f"  → 제목:{slide_title}, 본문:{len(body_texts)}, 표:{len(slide_tables)} 이미지: {len(slide_images)} 도형:{len(shape_texts)}, 링크:{len(slide_links)}")

    # ---------------------------------------------
    # 4) state 저장
    # ---------------------------------------------
    state["total_slides"] = slide_count

    state["titles"] = titles_list
    state["texts"] = texts_list
    state["tables"] = tables_list
    state["images"] = images_list
    state["slide_image"] = snapshots_list
    state["shape_texts"] = shape_texts_list
    state["links"] = links_list

    return state